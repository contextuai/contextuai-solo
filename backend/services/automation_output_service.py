"""
Automation Output Service

Processes output actions after an automation execution finishes.

Solo variant differences vs the enterprise reference:
- Files are written directly to ``~/.contextuai-solo/files/automations/``
  (the same root used by ``LocalStorageAdapter``). No ``FileStorageService``
  abstraction is required.
- PDFs are rendered with ``weasyprint``, PPTX with ``python-pptx``.
- Email uses ``smtplib`` (configurable via env vars) instead of AWS SES.
- The ``DISTRIBUTE`` action delegates to ``DistributionService.publish`` —
  reusing the existing LinkedIn/Twitter/blog/email/Slack outbound adapters.
"""

import csv
import io
import json
import logging
import os
import smtplib
import ssl
import uuid
from datetime import datetime
from email.mime.application import MIMEApplication
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from pathlib import Path
from typing import Any, Dict, List, Optional

from models.automation_models import (
    AutomationExecutionResponse,
    OutputAction,
    OutputActionType,
)

logger = logging.getLogger(__name__)


_FILES_DIR = Path(
    os.getenv(
        "AUTOMATIONS_OUTPUT_DIR",
        os.path.expanduser("~/.contextuai-solo/files/automations"),
    )
)


def _ensure_dir() -> Path:
    _FILES_DIR.mkdir(parents=True, exist_ok=True)
    return _FILES_DIR


def _safe_name(stem: str) -> str:
    keep = "".join(c if c.isalnum() or c in ("-", "_") else "_" for c in stem)
    return keep.strip("_") or "output"


def _timestamp() -> str:
    return datetime.utcnow().strftime("%Y%m%d_%H%M%S")


class AutomationOutputService:
    """Run the configured output actions for one automation execution."""

    async def process_output_actions(
        self,
        execution: AutomationExecutionResponse,
        output_actions: List[Dict[str, Any]],
        db=None,
    ) -> List[Dict[str, Any]]:
        results: List[Dict[str, Any]] = []
        for raw in output_actions:
            try:
                action = (
                    OutputAction(**raw)
                    if isinstance(raw, dict)
                    else raw
                )
                results.append(await self._dispatch(action, execution, db))
            except Exception as e:
                logger.exception("Output action failed")
                results.append(
                    {
                        "type": (raw.get("type") if isinstance(raw, dict) else str(raw)),
                        "status": "failed",
                        "error": str(e),
                    }
                )
        return results

    # ------------------------------------------------------------------
    # Dispatcher
    # ------------------------------------------------------------------

    async def _dispatch(
        self,
        action: OutputAction,
        execution: AutomationExecutionResponse,
        db,
    ) -> Dict[str, Any]:
        handlers = {
            OutputActionType.GENERATE_PDF: self._generate_pdf,
            OutputActionType.GENERATE_PPTX: self._generate_pptx,
            OutputActionType.SEND_EMAIL: self._send_email,
            OutputActionType.WEBHOOK: self._send_webhook,
            OutputActionType.SAVE_FILE: self._save_file,
            OutputActionType.DISTRIBUTE: self._distribute,
            OutputActionType.RUN_CODER_PROJECT: self._run_coder_project,
        }
        handler = handlers.get(action.type)
        if handler is None:
            return {
                "type": action.type.value,
                "status": "failed",
                "error": f"Unknown action type: {action.type}",
            }
        return await handler(execution, action.config, db)

    # ------------------------------------------------------------------
    # PDF / PPTX
    # ------------------------------------------------------------------

    async def _generate_pdf(
        self, execution: AutomationExecutionResponse, config: Dict[str, Any], db
    ) -> Dict[str, Any]:
        try:
            from weasyprint import HTML

            title = config.get("title") or "Automation Report"
            html = self._build_report_html(execution, title)
            pdf_bytes = HTML(string=html).write_pdf()

            outdir = _ensure_dir()
            filename = f"{_safe_name(title)}_{_timestamp()}.pdf"
            path = outdir / filename
            path.write_bytes(pdf_bytes)

            logger.info("PDF generated: %s (%d bytes)", path, len(pdf_bytes))
            return {
                "type": OutputActionType.GENERATE_PDF.value,
                "status": "success",
                "filename": filename,
                "path": str(path),
                "size": len(pdf_bytes),
            }
        except Exception as e:
            logger.exception("PDF generation failed")
            return {
                "type": OutputActionType.GENERATE_PDF.value,
                "status": "failed",
                "error": str(e),
            }

    async def _generate_pptx(
        self, execution: AutomationExecutionResponse, config: Dict[str, Any], db
    ) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            title = config.get("title") or "Automation Report"
            prs = Presentation()

            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            if slide.placeholders and len(slide.placeholders) > 1:
                slide.placeholders[1].text = (
                    f"Generated {execution.completed_at or execution.started_at}\n"
                    f"{execution.successful_steps}/{execution.total_steps} steps successful"
                )

            # One content slide per step
            for step in execution.steps:
                s = prs.slides.add_slide(prs.slide_layouts[1])
                s.shapes.title.text = f"Step {step.step_number}: @{step.persona}"
                body = s.placeholders[1].text_frame
                body.text = step.instruction or ""
                if step.result:
                    p = body.add_paragraph()
                    p.text = (step.result or "")[:1500]
                    p.font.size = Pt(12)

            # Summary slide
            if execution.final_result:
                s = prs.slides.add_slide(prs.slide_layouts[1])
                s.shapes.title.text = "Summary"
                s.placeholders[1].text_frame.text = execution.final_result[:3000]

            buf = io.BytesIO()
            prs.save(buf)
            pptx_bytes = buf.getvalue()

            outdir = _ensure_dir()
            filename = f"{_safe_name(title)}_{_timestamp()}.pptx"
            path = outdir / filename
            path.write_bytes(pptx_bytes)

            logger.info("PPTX generated: %s (%d bytes)", path, len(pptx_bytes))
            return {
                "type": OutputActionType.GENERATE_PPTX.value,
                "status": "success",
                "filename": filename,
                "path": str(path),
                "size": len(pptx_bytes),
            }
        except Exception as e:
            logger.exception("PPTX generation failed")
            return {
                "type": OutputActionType.GENERATE_PPTX.value,
                "status": "failed",
                "error": str(e),
            }

    # ------------------------------------------------------------------
    # Email (SMTP)
    # ------------------------------------------------------------------

    async def _send_email(
        self, execution: AutomationExecutionResponse, config: Dict[str, Any], db
    ) -> Dict[str, Any]:
        try:
            to_addrs = config.get("to") or []
            if isinstance(to_addrs, str):
                to_addrs = [to_addrs]
            if not to_addrs:
                return {
                    "type": OutputActionType.SEND_EMAIL.value,
                    "status": "failed",
                    "error": "No recipient addresses configured",
                }

            host = os.getenv("SMTP_HOST")
            if not host:
                return {
                    "type": OutputActionType.SEND_EMAIL.value,
                    "status": "failed",
                    "error": "SMTP not configured. Set SMTP_HOST/PORT/USER/PASSWORD or use the Distribute action with an Email connection.",
                }
            port = int(os.getenv("SMTP_PORT", "587"))
            user = os.getenv("SMTP_USER")
            password = os.getenv("SMTP_PASSWORD")
            from_addr = (
                config.get("from")
                or os.getenv("SMTP_FROM")
                or "noreply@contextuai-solo.local"
            )
            subject = config.get("subject") or f"Automation Report — {execution.automation_id}"

            msg = MIMEMultipart("alternative")
            msg["Subject"] = subject
            msg["From"] = from_addr
            msg["To"] = ", ".join(to_addrs)
            msg.attach(MIMEText(self._build_email_body(execution), "plain"))
            msg.attach(MIMEText(self._build_report_html(execution, subject), "html"))

            if config.get("include_pdf"):
                try:
                    from weasyprint import HTML
                    pdf_bytes = HTML(string=self._build_report_html(execution, subject)).write_pdf()
                    attach = MIMEApplication(pdf_bytes, _subtype="pdf")
                    attach.add_header(
                        "Content-Disposition",
                        "attachment",
                        filename=f"{_safe_name(subject)}.pdf",
                    )
                    msg.attach(attach)
                except Exception:
                    logger.warning("PDF attach failed; sending email without PDF", exc_info=True)

            ctx = ssl.create_default_context()
            with smtplib.SMTP(host, port, timeout=30) as smtp:
                smtp.ehlo()
                if port == 587:
                    smtp.starttls(context=ctx)
                    smtp.ehlo()
                if user and password:
                    smtp.login(user, password)
                smtp.sendmail(from_addr, to_addrs, msg.as_string())

            logger.info("Email sent to %d recipient(s) via %s", len(to_addrs), host)
            return {
                "type": OutputActionType.SEND_EMAIL.value,
                "status": "success",
                "details": {"recipients": len(to_addrs), "smtp_host": host},
            }
        except Exception as e:
            logger.exception("Email send failed")
            return {
                "type": OutputActionType.SEND_EMAIL.value,
                "status": "failed",
                "error": str(e),
            }

    # ------------------------------------------------------------------
    # Webhook
    # ------------------------------------------------------------------

    async def _send_webhook(
        self, execution: AutomationExecutionResponse, config: Dict[str, Any], db
    ) -> Dict[str, Any]:
        try:
            import httpx

            url = config.get("url")
            if not url:
                return {
                    "type": OutputActionType.WEBHOOK.value,
                    "status": "failed",
                    "error": "No webhook URL configured",
                }
            method = (config.get("method") or "POST").upper()
            headers = dict(config.get("headers") or {})
            headers.setdefault("Content-Type", "application/json")

            payload = {
                "execution_id": execution.execution_id,
                "automation_id": execution.automation_id,
                "status": execution.status.value if hasattr(execution.status, "value") else execution.status,
                "started_at": execution.started_at,
                "completed_at": execution.completed_at,
                "duration_ms": execution.duration_ms,
                "total_steps": execution.total_steps,
                "successful_steps": execution.successful_steps,
                "failed_steps": execution.failed_steps,
                "final_result": execution.final_result,
                "error_message": execution.error_message,
                "steps": [step.model_dump() for step in execution.steps],
            }

            async with httpx.AsyncClient(timeout=30.0) as client:
                if method == "PUT":
                    response = await client.put(url, json=payload, headers=headers)
                else:
                    response = await client.post(url, json=payload, headers=headers)

            return {
                "type": OutputActionType.WEBHOOK.value,
                "status": "success" if 200 <= response.status_code < 300 else "failed",
                "details": {"url": url, "status_code": response.status_code},
            }
        except Exception as e:
            logger.exception("Webhook failed")
            return {
                "type": OutputActionType.WEBHOOK.value,
                "status": "failed",
                "error": str(e),
            }

    # ------------------------------------------------------------------
    # Save raw file (json / csv / txt / md)
    # ------------------------------------------------------------------

    async def _save_file(
        self, execution: AutomationExecutionResponse, config: Dict[str, Any], db
    ) -> Dict[str, Any]:
        try:
            fmt = (config.get("format") or "json").lower()
            base = _safe_name(config.get("filename") or "automation_results")
            stamp = _timestamp()

            if fmt == "json":
                content = self._format_as_json(execution).encode("utf-8")
                filename = f"{base}_{stamp}.json"
            elif fmt == "csv":
                content = self._format_as_csv(execution).encode("utf-8")
                filename = f"{base}_{stamp}.csv"
            elif fmt in ("md", "markdown"):
                content = self._build_report_markdown(
                    execution, config.get("title") or "Automation Report"
                ).encode("utf-8")
                filename = f"{base}_{stamp}.md"
            else:
                content = self._build_email_body(execution).encode("utf-8")
                filename = f"{base}_{stamp}.txt"

            path = _ensure_dir() / filename
            path.write_bytes(content)
            return {
                "type": OutputActionType.SAVE_FILE.value,
                "status": "success",
                "filename": filename,
                "path": str(path),
                "size": len(content),
            }
        except Exception as e:
            logger.exception("Save-file failed")
            return {
                "type": OutputActionType.SAVE_FILE.value,
                "status": "failed",
                "error": str(e),
            }

    # ------------------------------------------------------------------
    # Distribute (publish through configured Distribution connection)
    # ------------------------------------------------------------------

    async def _distribute(
        self, execution: AutomationExecutionResponse, config: Dict[str, Any], db
    ) -> Dict[str, Any]:
        try:
            connection_id = config.get("connection_id") or config.get("channel_id")
            if not connection_id:
                return {
                    "type": OutputActionType.DISTRIBUTE.value,
                    "status": "failed",
                    "error": "connection_id is required for distribute action",
                }
            if db is None:
                return {
                    "type": OutputActionType.DISTRIBUTE.value,
                    "status": "failed",
                    "error": "Database handle not available — cannot publish",
                }

            from services.distribution_service import DistributionService

            svc = DistributionService(db)
            content = (
                config.get("content_template")
                or execution.final_result
                or self._build_email_body(execution)
            )
            title = config.get("title") or f"Automation: {execution.automation_id}"

            receipt = await svc.publish(
                channel_id=connection_id,
                content=content,
                title=title,
                published_by="automation",
            )
            return {
                "type": OutputActionType.DISTRIBUTE.value,
                "status": "success" if receipt.get("status") == "published" else "failed",
                "details": receipt,
            }
        except Exception as e:
            logger.exception("Distribute failed")
            return {
                "type": OutputActionType.DISTRIBUTE.value,
                "status": "failed",
                "error": str(e),
            }

    # ------------------------------------------------------------------
    # Run Coder project (headless)
    # ------------------------------------------------------------------

    async def _run_coder_project(
        self, execution: AutomationExecutionResponse, config: Dict[str, Any], db
    ) -> Dict[str, Any]:
        """Run a trusted Coder project headlessly and capture its output.

        Returns ``{type, status, project_id, output_lines, captured_output}``
        on success. ``captured_output`` is the joined stdout+stderr trimmed
        to 4000 chars so the artifact stays UI-friendly.
        """
        try:
            project_id = config.get("project_id")
            if not project_id:
                return {
                    "type": OutputActionType.RUN_CODER_PROJECT.value,
                    "status": "failed",
                    "error": "project_id is required",
                }
            timeout_seconds = int(config.get("timeout_seconds") or 60)

            if db is None:
                return {
                    "type": OutputActionType.RUN_CODER_PROJECT.value,
                    "status": "failed",
                    "error": "Database handle not available",
                }

            from repositories.coder_project_repository import CoderProjectRepository
            from services.coder_run_service import get_run_service

            repo = CoderProjectRepository(db)
            project = await repo.get_by_id(project_id)
            if not project:
                return {
                    "type": OutputActionType.RUN_CODER_PROJECT.value,
                    "status": "failed",
                    "project_id": project_id,
                    "error": f"Coder project '{project_id}' not found",
                }
            if not project.get("trusted"):
                return {
                    "type": OutputActionType.RUN_CODER_PROJECT.value,
                    "status": "failed",
                    "project_id": project_id,
                    "error": "Project is not trusted; cannot run headlessly.",
                }

            run_service = get_run_service()
            result = await run_service.run_headless(project, timeout_seconds)

            joined = "\n".join(result.get("output_lines") or [])
            captured = joined[:4000]

            return {
                "type": OutputActionType.RUN_CODER_PROJECT.value,
                "status": "success",
                "project_id": project_id,
                "output_lines": len(result.get("output_lines") or []),
                "captured_output": captured,
                "exit_code": result.get("exit_code"),
                "duration_seconds": result.get("duration_seconds"),
                "timed_out": result.get("timed_out", False),
            }
        except Exception as e:
            logger.exception("run_coder_project failed")
            return {
                "type": OutputActionType.RUN_CODER_PROJECT.value,
                "status": "failed",
                "error": str(e),
            }

    # ------------------------------------------------------------------
    # Formatters
    # ------------------------------------------------------------------

    def _build_report_markdown(
        self, execution: AutomationExecutionResponse, title: str
    ) -> str:
        status = execution.status.value if hasattr(execution.status, "value") else execution.status
        lines = [
            f"# {title}",
            "",
            f"- **Execution ID:** {execution.execution_id}",
            f"- **Status:** {status}",
            f"- **Started:** {execution.started_at}",
            f"- **Completed:** {execution.completed_at or 'N/A'}",
            f"- **Duration:** {execution.duration_ms or 0}ms",
            f"- **Steps:** {execution.successful_steps}/{execution.total_steps} successful",
            "",
        ]
        for step in execution.steps:
            step_status = step.status.value if hasattr(step.status, "value") else step.status
            lines.append(f"## Step {step.step_number} — @{step.persona}")
            lines.append(f"_Instruction:_ {step.instruction}")
            lines.append(f"_Status:_ `{step_status}` · {step.duration_ms}ms")
            lines.append("")
            lines.append(step.result or "_(no output)_")
            lines.append("")
        if execution.final_result:
            lines.append("## Summary")
            lines.append(execution.final_result)
        return "\n".join(lines)

    def _build_report_html(
        self, execution: AutomationExecutionResponse, title: str
    ) -> str:
        status = execution.status.value if hasattr(execution.status, "value") else execution.status
        rows: List[str] = []
        for step in execution.steps:
            step_status = step.status.value if hasattr(step.status, "value") else step.status
            color = (
                "#22c55e"
                if step_status == "success"
                else "#ef4444" if step_status == "failed" else "#f59e0b"
            )
            rows.append(
                f"""
                <tr>
                  <td style="padding:6px 8px;">{step.step_number}</td>
                  <td style="padding:6px 8px;">@{step.persona}</td>
                  <td style="padding:6px 8px;">{step.instruction}</td>
                  <td style="padding:6px 8px;color:{color};">{step_status}</td>
                  <td style="padding:6px 8px;">{step.duration_ms}ms</td>
                </tr>
                <tr><td colspan="5" style="padding:6px 8px 14px;color:#475569;font-size:12px;white-space:pre-wrap;">{(step.result or '')[:2000]}</td></tr>
                """
            )
        result_block = (
            f"""<div style="margin-top:20px;padding:12px;background:#f0fdf4;border-radius:8px;">
                <strong>Summary</strong>
                <div style="white-space:pre-wrap;margin-top:6px;">{execution.final_result}</div></div>"""
            if execution.final_result
            else ""
        )
        return f"""
        <html><head><meta charset="utf-8"></head>
        <body style="font-family:Arial,Helvetica,sans-serif;color:#0f172a;max-width:720px;margin:0 auto;padding:24px;">
          <h1 style="color:#FF4700;margin:0 0 12px;">{title}</h1>
          <div style="background:#f8fafc;padding:12px 16px;border-radius:8px;font-size:13px;">
            <div><strong>Status:</strong> {status}</div>
            <div><strong>Duration:</strong> {execution.duration_ms or 0}ms</div>
            <div><strong>Steps:</strong> {execution.successful_steps}/{execution.total_steps} successful</div>
          </div>
          <table style="width:100%;border-collapse:collapse;margin-top:16px;font-size:13px;">
            <thead><tr style="background:#e2e8f0;text-align:left;">
              <th style="padding:6px 8px;">#</th><th style="padding:6px 8px;">Agent</th>
              <th style="padding:6px 8px;">Instruction</th><th style="padding:6px 8px;">Status</th>
              <th style="padding:6px 8px;">Duration</th></tr></thead>
            <tbody>{''.join(rows)}</tbody>
          </table>
          {result_block}
          <p style="color:#94a3b8;font-size:11px;margin-top:24px;">Generated by ContextuAI Solo Automations</p>
        </body></html>
        """

    def _build_email_body(self, execution: AutomationExecutionResponse) -> str:
        status = execution.status.value if hasattr(execution.status, "value") else execution.status
        lines = [
            "Automation Execution Report",
            "=" * 40,
            f"Status: {status}",
            f"Duration: {execution.duration_ms or 0}ms",
            f"Steps: {execution.successful_steps}/{execution.total_steps} successful",
            "",
        ]
        for step in execution.steps:
            step_status = step.status.value if hasattr(step.status, "value") else step.status
            lines.append(f"Step {step.step_number} (@{step.persona}): {step_status}")
            lines.append(f"  {step.instruction}")
            preview = (step.result or "")[:240]
            if preview:
                lines.append(f"  → {preview}")
            lines.append("")
        if execution.final_result:
            lines.append("Final Result:")
            lines.append(execution.final_result)
        return "\n".join(lines)

    def _format_as_json(self, execution: AutomationExecutionResponse) -> str:
        return json.dumps(execution.model_dump(), indent=2, default=str)

    def _format_as_csv(self, execution: AutomationExecutionResponse) -> str:
        out = io.StringIO()
        w = csv.writer(out)
        w.writerow(["step", "persona", "instruction", "status", "duration_ms", "result"])
        for step in execution.steps:
            step_status = step.status.value if hasattr(step.status, "value") else step.status
            w.writerow(
                [
                    step.step_number,
                    step.persona,
                    step.instruction,
                    step_status,
                    step.duration_ms,
                    (step.result or "")[:500],
                ]
            )
        return out.getvalue()


automation_output_service = AutomationOutputService()

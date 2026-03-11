"""
Document Generation Service for AI Team Workspace.

Converts markdown content produced by AI agents into PDF, PPTX, and HTML formats
using Jinja2 templates and branded styling.

PPTX generation uses a two-step AI-powered process:
  1. Claude AI analyses markdown and produces structured slide JSON
  2. A professional renderer converts the JSON into styled PPTX slides
Falls back to basic regex parsing if the AI call is unavailable or fails.
"""

import os
import re
import json
import logging
import asyncio
from io import BytesIO
from datetime import datetime
from typing import Optional, List, Dict, Any

# Claude Agent SDK — used for AI-powered slide structuring
try:
    from claude_agent_sdk import (
        query as claude_query,
        ClaudeAgentOptions,
        AssistantMessage,
        ResultMessage,
        TextBlock,
    )
    CLAUDE_SDK_AVAILABLE = True
except ImportError:
    CLAUDE_SDK_AVAILABLE = False

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Claude prompt for structured slide generation
# ---------------------------------------------------------------------------

_SLIDE_GENERATION_PROMPT = """\
You are a presentation designer. Convert the following markdown document into a \
structured JSON array of slides for a professional PowerPoint presentation.

RULES:
- Output ONLY valid JSON — no markdown fences, no commentary.
- Aim for 10-20 slides total. Consolidate dense content; do not create a slide for every paragraph.
- Every bullet's "text" field must be plain text (no markdown formatting characters).
- Use "bold_ranges" to mark important words: each range is [start_char_index, end_char_index].
- Use these slide types:
  * "title"       — exactly one, first slide. Fields: title, subtitle.
  * "section"     — divider between major sections. Fields: title.
  * "content"     — standard bullet slide. Fields: title, bullets (array of {text, level (0 or 1), bold_ranges (optional)}).
  * "two_column"  — side-by-side comparison. Fields: title, left ({header, bullets}), right ({header, bullets}).
  * "table"       — tabular data. Fields: title, headers (array of strings), rows (array of arrays of strings).
  * "code"        — code snippet. Fields: title, language, code.
  * "summary"     — key takeaways (use at most once, near the end). Fields: title, bullets.
- Keep bullet text concise (under 100 chars per bullet).
- Prefer "content" slides for most material.
- Use "section" slides to separate major topics.
- Use "two_column" for comparisons, pros/cons, before/after.
- Use "table" only when the source has actual tabular data.
- Use "code" only when the source contains code blocks.

Respond with exactly: {"slides": [...]}

DOCUMENT TITLE: {title}

MARKDOWN CONTENT:
{content}
"""


class DocumentGenerationService:
    """Service for generating PDF, PPTX, and HTML documents from markdown content."""

    TEMPLATES_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates")

    def markdown_to_pdf(self, content: str, template: str = "report",
                        title: Optional[str] = None,
                        agent_count: Optional[int] = None) -> bytes:
        """Convert Markdown content to PDF using markdown2 + Jinja2 + weasyprint.

        Steps:
        1. Convert markdown to HTML using markdown2 (with extras like fenced-code-blocks, tables)
        2. Load Jinja2 template from templates/pdf/{template}.html
        3. Render HTML with template (inject converted markdown content)
        4. Load CSS from templates/pdf/styles.css
        5. Convert to PDF using weasyprint
        6. Return bytes

        If weasyprint is not available, fall back to returning the styled HTML as bytes.

        Args:
            content: Markdown string to convert.
            template: Name of the Jinja2 template (without extension) in templates/pdf/.
            title: Document title for the header. Defaults to "AI Workspace Report".
            agent_count: Number of agents that contributed to the document.

        Returns:
            PDF file content as bytes, or styled HTML bytes if weasyprint is unavailable.
        """
        if title is None:
            title = "AI Workspace Report"
        if agent_count is None:
            agent_count = 0

        date_str = datetime.utcnow().strftime("%B %d, %Y")

        # Step 1: Convert markdown to HTML
        try:
            import markdown2
            html_content = markdown2.markdown(
                content,
                extras=[
                    "fenced-code-blocks",
                    "tables",
                    "header-ids",
                    "code-friendly",
                    "cuddled-lists",
                    "strike",
                    "task_list",
                    "break-on-newline",
                ],
            )
        except ImportError:
            logger.warning("markdown2 not installed; using raw markdown wrapped in <pre> tags")
            html_content = f"<pre>{content}</pre>"

        # Step 2 & 3: Load Jinja2 template and render
        try:
            from jinja2 import Environment, FileSystemLoader

            pdf_templates_dir = os.path.join(self.TEMPLATES_PATH, "pdf")
            env = Environment(loader=FileSystemLoader(pdf_templates_dir))
            jinja_template = env.get_template(f"{template}.html")

            # Load CSS
            css_path = os.path.join(pdf_templates_dir, "styles.css")
            css_content = ""
            if os.path.exists(css_path):
                with open(css_path, "r", encoding="utf-8") as f:
                    css_content = f.read()

            rendered_html = jinja_template.render(
                content=html_content,
                title=title,
                date=date_str,
                agent_count=agent_count,
                css=css_content,
            )
        except ImportError:
            logger.warning("Jinja2 not installed; wrapping content in basic HTML")
            rendered_html = (
                f"<html><head><title>{title}</title></head>"
                f"<body><h1>{title}</h1><p>{date_str}</p>{html_content}</body></html>"
            )
        except Exception as e:
            logger.error("Failed to render Jinja2 template: %s", e)
            rendered_html = (
                f"<html><head><title>{title}</title></head>"
                f"<body><h1>{title}</h1><p>{date_str}</p>{html_content}</body></html>"
            )

        # Step 4 & 5: Convert to PDF using weasyprint
        try:
            from weasyprint import HTML as WeasyprintHTML

            pdf_bytes = WeasyprintHTML(string=rendered_html).write_pdf()
            logger.info("PDF generated successfully (%d bytes)", len(pdf_bytes))
            return pdf_bytes
        except ImportError:
            logger.warning(
                "weasyprint not installed; returning styled HTML as fallback. "
                "Install weasyprint for PDF generation: pip install weasyprint>=61.0"
            )
            return rendered_html.encode("utf-8")
        except Exception as e:
            logger.error("weasyprint PDF generation failed: %s", e)
            return rendered_html.encode("utf-8")

    # ------------------------------------------------------------------
    # PPTX generation — AI-powered slide structuring + professional renderer
    # ------------------------------------------------------------------

    def markdown_to_pptx(self, content: str, template: str = "slides",
                         title: Optional[str] = None) -> bytes:
        """Convert Markdown content to a professional PPTX presentation.

        Two-step process:
        1. Call Claude AI to convert markdown into structured slide JSON
           (slide types: title, section, content, two_column, table, code, summary).
        2. Render each slide using python-pptx with professional styling.

        Falls back to basic regex markdown parsing if the AI call fails.

        Args:
            content: Markdown string to convert.
            template: Template name (reserved for future use).
            title: Presentation title for the title slide.

        Returns:
            PPTX file content as bytes, or raw markdown bytes if python-pptx is unavailable.
        """
        if title is None:
            title = "AI Workspace Report"

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE

            import services.workspace.templates.pptx.slides_template as tpl
        except ImportError:
            logger.warning(
                "python-pptx not installed; returning raw markdown as fallback. "
                "Install python-pptx for PPTX generation: pip install python-pptx>=0.6.23"
            )
            return content.encode("utf-8")

        # ---- Try AI-powered slide generation, fall back to regex parsing -----

        slides_json: Optional[List[Dict[str, Any]]] = None

        if CLAUDE_SDK_AVAILABLE:
            try:
                slides_json = self._generate_slides_with_ai(content, title)
                if slides_json:
                    logger.info(
                        "AI slide generation succeeded: %d slides produced",
                        len(slides_json),
                    )
            except Exception as ai_err:
                logger.warning(
                    "AI slide generation failed, falling back to regex parsing: %s",
                    ai_err,
                )
                slides_json = None
        else:
            logger.info(
                "Claude Agent SDK not available; using regex-based slide parsing"
            )

        if slides_json is None:
            slides_json = self._parse_markdown_to_slides_legacy(content, title)

        # ---- Build the PPTX --------------------------------------------------

        prs = Presentation()
        prs.slide_width = Inches(tpl.SLIDE_WIDTH_INCHES)
        prs.slide_height = Inches(tpl.SLIDE_HEIGHT_INCHES)

        slide_number = 0
        for slide_data in slides_json:
            slide_number += 1
            try:
                self._render_slide(prs, slide_data, slide_number, tpl)
            except Exception as render_err:
                logger.error(
                    "Failed to render slide %d (type=%s): %s",
                    slide_number,
                    slide_data.get("type", "unknown"),
                    render_err,
                )
                # Render a simple fallback content slide so we don't lose material
                fallback = {
                    "type": "content",
                    "title": slide_data.get("title", f"Slide {slide_number}"),
                    "bullets": [{"text": str(slide_data), "level": 0}],
                }
                try:
                    self._render_slide(prs, fallback, slide_number, tpl)
                except Exception:
                    pass  # Skip entirely if even fallback fails

        # ---- Save to bytes ---------------------------------------------------

        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        pptx_bytes = buffer.read()

        logger.info(
            "PPTX generated successfully (%d bytes, %d slides)",
            len(pptx_bytes),
            slide_number,
        )
        return pptx_bytes

    # ------------------------------------------------------------------
    # Step 1: AI-powered slide JSON generation
    # ------------------------------------------------------------------

    def _generate_slides_with_ai(
        self, content: str, title: str
    ) -> Optional[List[Dict[str, Any]]]:
        """Call Claude to convert markdown into structured slide JSON.

        Returns a list of slide dicts on success, or None on failure.
        The caller is responsible for catching exceptions.
        """
        # Truncate very long content to avoid token limits (keep ~30k chars)
        max_content_chars = 30_000
        truncated = content[:max_content_chars]
        if len(content) > max_content_chars:
            truncated += "\n\n[Content truncated for slide generation]"

        prompt = _SLIDE_GENERATION_PROMPT.format(title=title, content=truncated)

        options = ClaudeAgentOptions(
            system_prompt=(
                "You are a JSON-only presentation structuring assistant. "
                "Output ONLY valid JSON. No explanations, no markdown fences."
            ),
            allowed_tools=[],
            permission_mode="default",
            model=os.getenv("CLAUDE_MODEL", "claude-sonnet-4-6"),
            max_turns=1,
        )

        # Run the async generator synchronously if we're not in an event loop,
        # or schedule it on the running loop.
        text_parts: List[str] = []

        async def _collect():
            async for message in claude_query(prompt=prompt, options=options):
                if isinstance(message, AssistantMessage):
                    for block in message.content:
                        if isinstance(block, TextBlock):
                            text_parts.append(block.text)

        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            loop = None

        if loop and loop.is_running():
            # We're inside an async context (e.g. Celery with async loop).
            # Create a new task and run until done via a helper thread.
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                pool.submit(asyncio.run, _collect()).result(timeout=120)
        else:
            asyncio.run(_collect())

        raw_text = "".join(text_parts).strip()
        if not raw_text:
            logger.warning("Claude returned empty response for slide generation")
            return None

        # Parse JSON — Claude may wrap in ```json ... ``` despite instructions
        json_text = raw_text
        fence_match = re.search(r"```(?:json)?\s*\n?(.*?)```", json_text, re.DOTALL)
        if fence_match:
            json_text = fence_match.group(1).strip()

        try:
            parsed = json.loads(json_text)
        except json.JSONDecodeError as e:
            logger.warning("Failed to parse AI slide JSON: %s\nRaw: %.500s", e, raw_text)
            return None

        slides = parsed.get("slides") if isinstance(parsed, dict) else parsed
        if not isinstance(slides, list) or len(slides) == 0:
            logger.warning("AI returned invalid slide structure: %.500s", raw_text)
            return None

        # Validate minimal structure
        valid_types = {"title", "section", "content", "two_column", "table", "code", "summary"}
        validated: List[Dict[str, Any]] = []
        for s in slides:
            if isinstance(s, dict) and s.get("type") in valid_types:
                validated.append(s)

        return validated if validated else None

    # ------------------------------------------------------------------
    # Legacy regex-based markdown parsing (fallback)
    # ------------------------------------------------------------------

    def _parse_markdown_to_slides_legacy(
        self, content: str, title: str
    ) -> List[Dict[str, Any]]:
        """Parse markdown via regex into basic slide JSON (legacy fallback).

        Returns a list of slide dicts compatible with the renderer.
        Always produces at least a title slide.
        """
        slides: List[Dict[str, Any]] = [
            {
                "type": "title",
                "title": title,
                "subtitle": f"Generated by ContextuAI AI Workspace - {datetime.utcnow().strftime('%B %d, %Y')}",
            }
        ]

        current_slide: Optional[Dict[str, Any]] = None

        for line in content.split("\n"):
            h_match = re.match(r"^#{1,2}\s+(.+)$", line)
            if h_match:
                if current_slide is not None:
                    slides.append(current_slide)
                current_slide = {
                    "type": "content",
                    "title": h_match.group(1).strip(),
                    "bullets": [],
                }
                continue

            if current_slide is not None:
                stripped = line.strip()
                if stripped:
                    # Determine indent level
                    indent_chars = len(line) - len(line.lstrip())
                    level = 1 if indent_chars >= 4 or line.startswith("  ") else 0
                    text = re.sub(r"^[\-\*\+]\s+", "", stripped)
                    text = re.sub(r"^\d+\.\s+", "", text)
                    current_slide["bullets"].append({"text": text, "level": level})
            else:
                if line.strip():
                    current_slide = {
                        "type": "content",
                        "title": title,
                        "bullets": [{"text": line.strip(), "level": 0}],
                    }

        if current_slide is not None:
            slides.append(current_slide)

        return slides

    # ------------------------------------------------------------------
    # Step 2: Professional PPTX renderer
    # ------------------------------------------------------------------

    @staticmethod
    def _hex_to_rgb(hex_str: str):
        """Convert a 6-char hex string to pptx RGBColor."""
        from pptx.dml.color import RGBColor
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

    def _render_slide(self, prs, slide_data: Dict[str, Any], slide_number: int, tpl) -> None:
        """Dispatch to the appropriate slide renderer based on slide type."""
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        slide_type = slide_data.get("type", "content")

        renderer_map = {
            "title": self._render_title_slide,
            "section": self._render_section_slide,
            "content": self._render_content_slide,
            "two_column": self._render_two_column_slide,
            "table": self._render_table_slide,
            "code": self._render_code_slide,
            "summary": self._render_summary_slide,
        }

        renderer = renderer_map.get(slide_type, self._render_content_slide)
        renderer(prs, slide_data, tpl)

        # Add slide number footer and watermark to every slide
        slide = prs.slides[len(prs.slides) - 1]
        self._add_slide_footer(slide, slide_number, len(prs.slides), tpl)

    # ---- Title slide ----

    def _render_title_slide(self, prs, data: Dict[str, Any], tpl) -> None:
        """Render the title slide with brand gradient background."""
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.oxml.ns import qn
        import lxml.etree as etree

        # Use blank layout so we have full control
        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # Set gradient background
        bg = slide.background
        fill = bg.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR)
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[1].color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR_DARK)
        fill.gradient_stops[1].position = 1.0

        slide_w = Inches(tpl.SLIDE_WIDTH_INCHES)
        slide_h = Inches(tpl.SLIDE_HEIGHT_INCHES)

        # Title text box — centered
        title_text = data.get("title", "Presentation")
        left = Inches(1.5)
        top = Inches(2.0)
        width = slide_w - Inches(3.0)
        height = Inches(2.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = tpl.TITLE_FONT
        run.font.size = Pt(44)
        run.font.color.rgb = self._hex_to_rgb(tpl.WHITE)
        run.font.bold = True

        # Subtitle
        subtitle_text = data.get(
            "subtitle",
            f"Generated by ContextuAI AI Workspace - {datetime.utcnow().strftime('%B %d, %Y')}",
        )
        sub_top = Inches(4.2)
        sub_height = Inches(1.2)
        subBox = slide.shapes.add_textbox(left, sub_top, width, sub_height)
        stf = subBox.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = subtitle_text
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.runs[0]
        sr.font.name = tpl.BODY_FONT
        sr.font.size = Pt(22)
        sr.font.color.rgb = self._hex_to_rgb(tpl.WHITE)

        # Decorative accent line at bottom
        from pptx.enum.shapes import MSO_SHAPE
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.0), Inches(5.8), Inches(5.333), Inches(0.05),
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = self._hex_to_rgb(tpl.WHITE)
        line_shape.line.fill.background()

    # ---- Section divider slide ----

    def _render_section_slide(self, prs, data: Dict[str, Any], tpl) -> None:
        """Render a section divider with dark background."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE

        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # Dark background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(tpl.SECTION_BG)

        slide_w = Inches(tpl.SLIDE_WIDTH_INCHES)

        # Accent bar on left
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(2.5), Inches(0.12), Inches(2.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR)
        bar.line.fill.background()

        # Title
        title_text = data.get("title", "Section")
        txBox = slide.shapes.add_textbox(
            Inches(1.3), Inches(2.6), slide_w - Inches(2.5), Inches(2.2),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        run = p.runs[0]
        run.font.name = tpl.TITLE_FONT
        run.font.size = Pt(38)
        run.font.color.rgb = self._hex_to_rgb(tpl.WHITE)
        run.font.bold = True

    # ---- Content slide ----

    def _render_content_slide(self, prs, data: Dict[str, Any], tpl) -> None:
        """Render a content slide with title bar strip and styled bullets."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # White background (default)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(tpl.WHITE)

        slide_w = Inches(tpl.SLIDE_WIDTH_INCHES)
        cfg = tpl.SLIDE_LAYOUTS.get("content", {})

        # Orange accent strip at top
        strip_height = Inches(cfg.get("title_bar_height_inches", 0.06))
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), slide_w, strip_height,
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR)
        strip.line.fill.background()

        # Title
        title_text = data.get("title", "")
        title_box = slide.shapes.add_textbox(
            Inches(tpl.MARGIN_LEFT), Inches(0.4),
            slide_w - Inches(tpl.MARGIN_LEFT + tpl.MARGIN_RIGHT), Inches(1.0),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        run = p.runs[0]
        run.font.name = cfg.get("title_font", tpl.TITLE_FONT)
        run.font.size = Pt(cfg.get("title_font_size", 28))
        run.font.color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR)
        run.font.bold = True

        # Bullets
        bullets = data.get("bullets", [])
        if not bullets:
            return

        body_box = slide.shapes.add_textbox(
            Inches(tpl.MARGIN_LEFT), Inches(tpl.CONTENT_TOP),
            slide_w - Inches(tpl.MARGIN_LEFT + tpl.MARGIN_RIGHT),
            Inches(tpl.SLIDE_HEIGHT_INCHES - tpl.CONTENT_TOP - tpl.MARGIN_BOTTOM),
        )
        btf = body_box.text_frame
        btf.word_wrap = True

        for idx, bullet in enumerate(bullets):
            if isinstance(bullet, str):
                bullet = {"text": bullet, "level": 0}

            text = bullet.get("text", "")
            level = bullet.get("level", 0)
            bold_ranges = bullet.get("bold_ranges", [])

            if idx == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()

            bp.level = level
            bp.space_after = Pt(6)

            # Handle bold ranges by splitting text into runs
            if bold_ranges:
                self._add_formatted_runs(
                    bp, text, bold_ranges,
                    font_name=cfg.get("body_font", tpl.BODY_FONT),
                    font_size=Pt(cfg.get("body_font_size", 18)),
                    font_color=self._hex_to_rgb(tpl.DARK_GRAY),
                )
            else:
                bp.text = text
                for r in bp.runs:
                    r.font.name = cfg.get("body_font", tpl.BODY_FONT)
                    r.font.size = Pt(cfg.get("body_font_size", 18))
                    r.font.color.rgb = self._hex_to_rgb(tpl.DARK_GRAY)

            # Add bullet marker color for level-0 items
            if level == 0 and bp.runs:
                # Bullet color is managed by the paragraph run font
                pass

    # ---- Two-column slide ----

    def _render_two_column_slide(self, prs, data: Dict[str, Any], tpl) -> None:
        """Render a two-column comparison slide."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(tpl.WHITE)

        slide_w = Inches(tpl.SLIDE_WIDTH_INCHES)
        cfg = tpl.SLIDE_LAYOUTS.get("two_column", {})

        # Orange accent strip
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), slide_w, Inches(0.06),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR)
        strip.line.fill.background()

        # Title
        title_text = data.get("title", "")
        title_box = slide.shapes.add_textbox(
            Inches(tpl.MARGIN_LEFT), Inches(0.4),
            slide_w - Inches(tpl.MARGIN_LEFT + tpl.MARGIN_RIGHT), Inches(1.0),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        run = p.runs[0]
        run.font.name = cfg.get("title_font", tpl.TITLE_FONT)
        run.font.size = Pt(cfg.get("title_font_size", 28))
        run.font.color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR)
        run.font.bold = True

        # Column geometry
        gap = cfg.get("column_gap_inches", 0.5)
        usable_w = tpl.SLIDE_WIDTH_INCHES - tpl.MARGIN_LEFT - tpl.MARGIN_RIGHT - gap
        col_w = usable_w / 2
        col_top = tpl.CONTENT_TOP
        col_height = tpl.SLIDE_HEIGHT_INCHES - col_top - tpl.MARGIN_BOTTOM

        left_data = data.get("left", {})
        right_data = data.get("right", {})

        for col_idx, col_data in enumerate([left_data, right_data]):
            x = tpl.MARGIN_LEFT + col_idx * (col_w + gap)

            # Column header
            header = col_data.get("header", "")
            hdr_box = slide.shapes.add_textbox(
                Inches(x), Inches(col_top - 0.5), Inches(col_w), Inches(0.5),
            )
            htf = hdr_box.text_frame
            htf.word_wrap = True
            hp = htf.paragraphs[0]
            hp.text = header
            hr = hp.runs[0]
            hr.font.name = cfg.get("title_font", tpl.TITLE_FONT)
            hr.font.size = Pt(cfg.get("column_header_font_size", 20))
            hr.font.color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR)
            hr.font.bold = True

            # Separator line under header
            sep = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(col_top - 0.02), Inches(col_w), Inches(0.03),
            )
            sep.fill.solid()
            sep.fill.fore_color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR_LIGHT)
            sep.line.fill.background()

            # Bullets
            bullets = col_data.get("bullets", [])
            if bullets:
                bx = slide.shapes.add_textbox(
                    Inches(x), Inches(col_top + 0.15), Inches(col_w), Inches(col_height - 0.15),
                )
                btf = bx.text_frame
                btf.word_wrap = True

                for b_idx, bullet in enumerate(bullets):
                    if isinstance(bullet, str):
                        bullet = {"text": bullet, "level": 0}
                    if b_idx == 0:
                        bp = btf.paragraphs[0]
                    else:
                        bp = btf.add_paragraph()
                    bp.text = bullet.get("text", "")
                    bp.level = bullet.get("level", 0)
                    bp.space_after = Pt(4)
                    for r in bp.runs:
                        r.font.name = cfg.get("body_font", tpl.BODY_FONT)
                        r.font.size = Pt(cfg.get("body_font_size", 16))
                        r.font.color.rgb = self._hex_to_rgb(tpl.DARK_GRAY)

    # ---- Table slide ----

    def _render_table_slide(self, prs, data: Dict[str, Any], tpl) -> None:
        """Render a table slide with styled header row."""
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN

        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(tpl.WHITE)

        slide_w = Inches(tpl.SLIDE_WIDTH_INCHES)
        cfg = tpl.SLIDE_LAYOUTS.get("table", {})

        # Title
        title_text = data.get("title", "")
        title_box = slide.shapes.add_textbox(
            Inches(tpl.MARGIN_LEFT), Inches(0.4),
            slide_w - Inches(tpl.MARGIN_LEFT + tpl.MARGIN_RIGHT), Inches(1.0),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        run = p.runs[0]
        run.font.name = cfg.get("title_font", tpl.TITLE_FONT)
        run.font.size = Pt(cfg.get("title_font_size", 28))
        run.font.color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR)
        run.font.bold = True

        headers = data.get("headers", [])
        rows = data.get("rows", [])

        if not headers:
            return

        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header row

        table_width = tpl.SLIDE_WIDTH_INCHES - tpl.MARGIN_LEFT - tpl.MARGIN_RIGHT
        table_height = min(
            tpl.SLIDE_HEIGHT_INCHES - tpl.CONTENT_TOP - tpl.MARGIN_BOTTOM,
            0.4 * num_rows + 0.1,
        )

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(tpl.MARGIN_LEFT), Inches(tpl.CONTENT_TOP),
            Inches(table_width), Inches(table_height),
        )
        table = table_shape.table

        # Style header row
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header_text)
            # Header background
            cell_fill = cell.fill
            cell_fill.solid()
            cell_fill.fore_color.rgb = self._hex_to_rgb(tpl.TABLE_HEADER_BG)
            # Header text
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for r in para.runs:
                    r.font.name = tpl.BODY_FONT
                    r.font.size = Pt(cfg.get("body_font_size", 14))
                    r.font.color.rgb = self._hex_to_rgb(tpl.TABLE_HEADER_TEXT)
                    r.font.bold = True

        # Style data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx in range(num_cols):
                cell = table.cell(row_idx + 1, col_idx)
                cell_val = str(row_data[col_idx]) if col_idx < len(row_data) else ""
                cell.text = cell_val

                # Alternating row background
                if row_idx % 2 == 1:
                    cell_fill = cell.fill
                    cell_fill.solid()
                    cell_fill.fore_color.rgb = self._hex_to_rgb(tpl.TABLE_ROW_ALT)

                for para in cell.text_frame.paragraphs:
                    for r in para.runs:
                        r.font.name = tpl.BODY_FONT
                        r.font.size = Pt(cfg.get("body_font_size", 14))
                        r.font.color.rgb = self._hex_to_rgb(tpl.DARK_GRAY)

    # ---- Code slide ----

    def _render_code_slide(self, prs, data: Dict[str, Any], tpl) -> None:
        """Render a code slide with monospace font and gray background."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(tpl.WHITE)

        slide_w = Inches(tpl.SLIDE_WIDTH_INCHES)
        cfg = tpl.SLIDE_LAYOUTS.get("code", {})

        # Title
        title_text = data.get("title", "Code")
        title_box = slide.shapes.add_textbox(
            Inches(tpl.MARGIN_LEFT), Inches(0.4),
            slide_w - Inches(tpl.MARGIN_LEFT + tpl.MARGIN_RIGHT), Inches(1.0),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        run = p.runs[0]
        run.font.name = cfg.get("title_font", tpl.TITLE_FONT)
        run.font.size = Pt(cfg.get("title_font_size", 28))
        run.font.color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR)
        run.font.bold = True

        # Language label
        language = data.get("language", "")
        if language:
            lang_box = slide.shapes.add_textbox(
                Inches(tpl.MARGIN_LEFT), Inches(1.45),
                Inches(3.0), Inches(0.3),
            )
            ltf = lang_box.text_frame
            lp = ltf.paragraphs[0]
            lp.text = language.upper()
            lr = lp.runs[0]
            lr.font.name = tpl.BODY_FONT
            lr.font.size = Pt(cfg.get("language_label_font_size", 10))
            lr.font.color.rgb = self._hex_to_rgb(tpl.MEDIUM_GRAY)
            lr.font.bold = True

        # Code background rectangle
        code_left = Inches(tpl.MARGIN_LEFT)
        code_top = Inches(tpl.CONTENT_TOP)
        code_width = slide_w - Inches(tpl.MARGIN_LEFT + tpl.MARGIN_RIGHT)
        code_height = Inches(tpl.SLIDE_HEIGHT_INCHES - tpl.CONTENT_TOP - tpl.MARGIN_BOTTOM)

        code_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            code_left, code_top, code_width, code_height,
        )
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = self._hex_to_rgb(cfg.get("code_bg_color", tpl.CODE_BG))
        code_bg.line.color.rgb = self._hex_to_rgb(cfg.get("code_border_color", tpl.CODE_BORDER))
        code_bg.line.width = Pt(1)

        # Code text on top of the background
        code_text = data.get("code", "")
        # Truncate very long code blocks to fit on a slide
        lines = code_text.split("\n")
        if len(lines) > 25:
            lines = lines[:24] + ["  // ... (truncated)"]
        code_text = "\n".join(lines)

        code_box = slide.shapes.add_textbox(
            code_left + Inches(0.3), code_top + Inches(0.2),
            code_width - Inches(0.6), code_height - Inches(0.4),
        )
        ctf = code_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = code_text
        cr = cp.runs[0]
        cr.font.name = cfg.get("code_font", tpl.CODE_FONT)
        cr.font.size = Pt(cfg.get("code_font_size", 12))
        cr.font.color.rgb = self._hex_to_rgb(cfg.get("code_color", tpl.DARK_GRAY))

    # ---- Summary slide ----

    def _render_summary_slide(self, prs, data: Dict[str, Any], tpl) -> None:
        """Render a summary slide with numbered takeaways and highlight background."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        cfg = tpl.SLIDE_LAYOUTS.get("summary", {})

        # Warm highlight background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(cfg.get("background_color", tpl.SUMMARY_BG))

        slide_w = Inches(tpl.SLIDE_WIDTH_INCHES)

        # Orange accent strip at top
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), slide_w, Inches(0.08),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR)
        strip.line.fill.background()

        # Title
        title_text = data.get("title", "Key Takeaways")
        title_box = slide.shapes.add_textbox(
            Inches(tpl.MARGIN_LEFT), Inches(0.5),
            slide_w - Inches(tpl.MARGIN_LEFT + tpl.MARGIN_RIGHT), Inches(1.0),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        run = p.runs[0]
        run.font.name = cfg.get("title_font", tpl.TITLE_FONT)
        run.font.size = Pt(cfg.get("title_font_size", 32))
        run.font.color.rgb = self._hex_to_rgb(tpl.BRAND_COLOR)
        run.font.bold = True

        # Numbered takeaways
        bullets = data.get("bullets", [])
        if not bullets:
            return

        body_box = slide.shapes.add_textbox(
            Inches(tpl.MARGIN_LEFT + 0.2), Inches(tpl.CONTENT_TOP),
            slide_w - Inches(tpl.MARGIN_LEFT + tpl.MARGIN_RIGHT + 0.4),
            Inches(tpl.SLIDE_HEIGHT_INCHES - tpl.CONTENT_TOP - tpl.MARGIN_BOTTOM),
        )
        btf = body_box.text_frame
        btf.word_wrap = True

        for idx, bullet in enumerate(bullets):
            if isinstance(bullet, str):
                bullet = {"text": bullet, "level": 0}

            text = bullet.get("text", "")
            numbered_text = f"{idx + 1}.  {text}"

            if idx == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()

            bp.space_after = Pt(10)
            bp.space_before = Pt(4)

            # Add numbered run with colored number
            from pptx.oxml.ns import qn

            num_run = bp.add_run()
            num_run.text = f"{idx + 1}. "
            num_run.font.name = cfg.get("body_font", tpl.BODY_FONT)
            num_run.font.size = Pt(cfg.get("number_font_size", 22))
            num_run.font.color.rgb = self._hex_to_rgb(cfg.get("number_color", tpl.BRAND_COLOR))
            num_run.font.bold = True

            text_run = bp.add_run()
            text_run.text = text
            text_run.font.name = cfg.get("body_font", tpl.BODY_FONT)
            text_run.font.size = Pt(cfg.get("body_font_size", 18))
            text_run.font.color.rgb = self._hex_to_rgb(tpl.DARK_GRAY)

    # ---- Footer and branding ----

    def _add_slide_footer(self, slide, slide_number: int, total_slides: int, tpl) -> None:
        """Add slide number and ContextuAI watermark to a slide."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide_w = Inches(tpl.SLIDE_WIDTH_INCHES)
        slide_h = Inches(tpl.SLIDE_HEIGHT_INCHES)

        # Slide number (bottom right)
        num_box = slide.shapes.add_textbox(
            slide_w - Inches(1.5), slide_h - Inches(0.4),
            Inches(1.2), Inches(0.3),
        )
        ntf = num_box.text_frame
        np_ = ntf.paragraphs[0]
        np_.text = str(slide_number)
        np_.alignment = PP_ALIGN.RIGHT
        nr = np_.runs[0]
        nr.font.name = tpl.BODY_FONT
        nr.font.size = Pt(tpl.FOOTER_FONT_SIZE)
        nr.font.color.rgb = self._hex_to_rgb(tpl.FOOTER_COLOR)

        # Watermark text (bottom left)
        wm_box = slide.shapes.add_textbox(
            Inches(0.4), slide_h - Inches(0.4),
            Inches(3.0), Inches(0.3),
        )
        wtf = wm_box.text_frame
        wp = wtf.paragraphs[0]
        wp.text = tpl.FOOTER_TEXT
        wp.alignment = PP_ALIGN.LEFT
        wr = wp.runs[0]
        wr.font.name = tpl.BODY_FONT
        wr.font.size = Pt(tpl.FOOTER_FONT_SIZE)
        wr.font.color.rgb = self._hex_to_rgb(tpl.WATERMARK_COLOR)

    # ---- Shared formatting helpers ----

    def _add_formatted_runs(
        self, paragraph, text: str, bold_ranges: List[List[int]],
        font_name: str, font_size, font_color,
    ) -> None:
        """Add runs to a paragraph with bold formatting for specified character ranges."""
        if not bold_ranges:
            paragraph.text = text
            for r in paragraph.runs:
                r.font.name = font_name
                r.font.size = font_size
                r.font.color.rgb = font_color
            return

        # Sort ranges and build run segments
        bold_ranges = sorted(bold_ranges, key=lambda x: x[0])
        pos = 0
        for br in bold_ranges:
            start = max(br[0], pos)
            end = min(br[1], len(text))
            if start > pos:
                # Normal text before bold
                r = paragraph.add_run()
                r.text = text[pos:start]
                r.font.name = font_name
                r.font.size = font_size
                r.font.color.rgb = font_color
            if start < end:
                # Bold text
                r = paragraph.add_run()
                r.text = text[start:end]
                r.font.name = font_name
                r.font.size = font_size
                r.font.color.rgb = font_color
                r.font.bold = True
            pos = end

        if pos < len(text):
            r = paragraph.add_run()
            r.text = text[pos:]
            r.font.name = font_name
            r.font.size = font_size
            r.font.color.rgb = font_color

    def markdown_to_html(self, content: str, template: str = "report",
                         title: Optional[str] = None,
                         agent_count: Optional[int] = None) -> bytes:
        """Convert Markdown content to a standalone HTML document.

        Uses markdown2 for conversion with the same extras as PDF generation,
        then wraps the result in a branded HTML page with inline CSS.

        Args:
            content: Markdown string to convert.
            template: Template name (report, brief).
            title: Document title.
            agent_count: Number of contributing agents.

        Returns:
            HTML file content as UTF-8 bytes.
        """
        if title is None:
            title = "AI Workspace Report"
        if agent_count is None:
            agent_count = 0

        date_str = datetime.utcnow().strftime("%B %d, %Y")

        # Convert markdown to HTML
        try:
            import markdown2
            html_content = markdown2.markdown(
                content,
                extras=[
                    "fenced-code-blocks",
                    "tables",
                    "header-ids",
                    "code-friendly",
                    "cuddled-lists",
                    "strike",
                    "task_list",
                    "break-on-newline",
                ],
            )
        except ImportError:
            logger.warning("markdown2 not installed; using raw markdown in <pre> tags")
            html_content = f"<pre>{content}</pre>"

        # Try to use Jinja2 template (same as PDF pipeline)
        try:
            from jinja2 import Environment, FileSystemLoader

            pdf_templates_dir = os.path.join(self.TEMPLATES_PATH, "pdf")
            env = Environment(loader=FileSystemLoader(pdf_templates_dir))
            jinja_template = env.get_template(f"{template}.html")

            css_path = os.path.join(pdf_templates_dir, "styles.css")
            css_content = ""
            if os.path.exists(css_path):
                with open(css_path, "r", encoding="utf-8") as f:
                    css_content = f.read()

            rendered_html = jinja_template.render(
                content=html_content,
                title=title,
                date=date_str,
                agent_count=agent_count,
                css=css_content,
            )
        except Exception as e:
            logger.warning("Jinja2 template rendering failed for HTML export: %s", e)
            rendered_html = (
                f"<!DOCTYPE html><html><head><meta charset='utf-8'>"
                f"<title>{title}</title>"
                f"<style>body{{font-family:system-ui,sans-serif;max-width:800px;margin:0 auto;padding:2rem;line-height:1.6}}"
                f"h1,h2,h3{{color:#FF4700}}table{{border-collapse:collapse;width:100%}}"
                f"th,td{{border:1px solid #ddd;padding:8px;text-align:left}}"
                f"code{{background:#f4f4f4;padding:2px 6px;border-radius:3px}}"
                f"pre{{background:#f4f4f4;padding:1rem;border-radius:6px;overflow-x:auto}}</style>"
                f"</head><body><h1>{title}</h1><p><em>{date_str}</em></p>"
                f"{html_content}</body></html>"
            )

        logger.info("HTML document generated successfully (%d bytes)", len(rendered_html))
        return rendered_html.encode("utf-8")
